from pptx import Presentation
from pptx.util import Mm, Pt
from pptx.dml.color import RGBColor
from Text import TextSong
from Setting import TextType
from Setting import TextColor
from Setting import Pattern
from Setting import PPT_SONG
import os
import re


class PPTCreator:
    def __init__(self, back_color, sample_file_name=None) -> None:
        self.prs = Presentation(sample_file_name)
        self.set_slide_layout()
        if sample_file_name is None:
            self.set_slide_back_color(back_color)

    def generate_PPT(self, file_name):
        # 프레젠테이션 파일 저장
        desktop_directory = os.path.join(os.path.expanduser("~"), "Desktop")
        self.prs.save(f"{desktop_directory}/{file_name}.pptx")

    def create_slide(self, Text_list):
        slides = []
        slide = self.add_new_slide()
        while Text_list:
            Text = Text_list.pop(0)
            if Text.get_text_type() == TextType.FILE_NAME:
                self.join_text(slide, Text)
                slides.append(slide)
                file_name = str(Text)
                slide = self.add_new_slide()
            if (
                Text.get_text_type() == TextType.MENT_GUIDE
                or Text.get_text_type() == TextType.MENT_GUIDE_INTRO
            ):
                self.join_text(slide, Text)
                self.enter(slide)
                # if "없음" in str(Text):
                # self.enter(slide)
            if Text.get_text_type() == TextType.MENT:
                self.join_text(slide, Text)
                slides.append(slide)
                slide = self.add_new_slide()
            if Text.get_text_type() == TextType.SONG_TITLE:
                self.join_text(slide, Text)
                self.enter_new_line(slide)
            if Text.get_text_type() == TextType.LYRICS:
                self.join_text(slide, Text)
                slides.append(slide)
                slide = self.add_new_slide()
            if Text.get_text_type() == TextType.LYRICS_GUIDE:
                self.join_text(slide, Text)
            if (
                Text.get_text_type() == TextType.INTERLUDE
                or Text.get_text_type() == TextType.INTRO
                or Text.get_text_type() == TextType.ELSE
            ):
                self.join_text(slide, Text)
                self.enter(slide)
        slides.append(slide)
        previews = self.get_previews(self.prs)
        self.add_previews(slides, previews)
        self.slide_end(slides[-1])

        self.generate_PPT(file_name)

    def set_slide_layout(self):
        # 슬라이드 크기를 16:9 비율(와이드스크린)로 조절
        self.prs.slide_width = Mm(338.67)
        self.prs.slide_height = Mm(190.5)

        # 제목만 있는 슬라이드 레이아웃 가져오기
        self.slide_layout = self.prs.slide_master.slide_layouts[5]

    # 배경색 설정
    def set_slide_back_color(self, color):
        self.prs.slide_master.background.fill.solid()
        self.prs.slide_master.background.fill.fore_color.rgb = color

    def add_new_slide(self):
        # 슬라이드 추가
        slide = self.prs.slides.add_slide(self.slide_layout)

        # 제목 추가
        title_shape = slide.shapes.title

        # 텍스트 상자의 크기를 밀리미터 단위로 설정
        title_shape.width = Mm(338.67)  # 가로 길이를 338.67밀리미터로 설정
        title_shape.height = Mm(50)  # 세로 길이를 50밀리미터로 설정

        # 제목 텍스트 위치 설정 (가운데)
        title_shape.left = int((self.prs.slide_width - title_shape.width) / 2)
        title_shape.top = int((self.prs.slide_height - title_shape.height) / 2)

        return slide

    def enter_new_line(self, slide):
        title_shape = slide.shapes.title
        title_text_frame = title_shape.text_frame
        title_text_frame.add_paragraph()
        title_text_frame.add_paragraph()

    def enter(self, slide):
        title_shape = slide.shapes.title
        title_text_frame = title_shape.text_frame
        title_text_frame.add_paragraph()

    def join_text(self, slide, text):
        title_shape = slide.shapes.title
        title_text_frame = title_shape.text_frame
        p = title_text_frame.paragraphs[-1]  # 마지막 단락 선택
        run = p.add_run()
        # 분류 다 된 후 Text 클래스로 전달 받을 때
        if isinstance(text, TextSong):
            self.set_text(
                run, str(text), PPT_SONG.font, Pt(PPT_SONG.size), text.get_text_color()
            )
        # 슬라이드 마지막 줄에 미리보기 추가할 때
        else:
            self.set_text(run, text, PPT_SONG.font, Pt(PPT_SONG.size), TextColor.YELLOW)

    def join_text_end(self, slide, text):
        title_shape = slide.shapes.title
        title_text_frame = title_shape.text_frame
        p = title_text_frame.paragraphs[-1]  # 마지막 단락 선택
        run = p.add_run()
        if "마무리" in text or "끝" in text:
            self.set_text(run, text, PPT_SONG.font, Pt(PPT_SONG.size), TextColor.RED)
        else:
            self.set_text(run, text, PPT_SONG.font, Pt(PPT_SONG.size), TextColor.YELLOW)

    def get_previews(self, prs):
        previews = []
        for slide in prs.slides:
            title_shape = slide.shapes.title
            text = title_shape.text_frame.text
            first_line = text.split("\n")[0]

            # 곡목이면
            if bool(re.search(Pattern.song_title, first_line)):
                previews.append("- " + first_line.strip() + " -")
            elif "멘트" in first_line:
                if "없음" in first_line:
                    previews.append("- 멘트 x -")
                else:
                    previews.append("- 멘트 o -")
                if "마무리" in first_line:
                    previews.append("- 마무리" + previews.pop()[1:])
                elif "간주" in first_line:
                    previews.append("- 간주" + previews.pop()[1:])
            elif first_line == "":
                previews.append("- 끝 -")
            else:
                previews.append("- " + first_line.strip() + " -")
        return previews

    def add_previews(self, slides, previews):
        for i in range(1, len(slides) - 1):
            self.enter_new_line(slides[i])

            # 마지막 프리뷰 빨간 글자. for문 범위가 -1 이전이니 -2가 마지막
            if i == len(slides) - 2:
                self.join_text_end(slides[i], previews[i + 1])
                return

            self.join_text(slides[i], previews[i + 1])

            # 멘트 없는 경우에  프리뷰로만 표시하고 다음 슬라이드에는 필요 없음.
            if "x" in previews[i + 1]:
                slides[i + 1].shapes.title.text_frame.paragraphs[0].text = ""
                slides[i + 1].shapes.title.text_frame.paragraphs[0].font.size = Pt(1)

    def slide_end(self, slide):
        title_shape = slide.shapes.title
        p = title_shape.text_frame.paragraphs[0]  # 첫번째 단락 선택. 단락? 줄?
        if "마무리" in p.text:
            for run in p.runs:
                run.font.color.rgb = TextColor.RED

    # 폰트와 글자 크기가 안 바뀐다면 일일이 변수로 받아서 설정할 필요 없음.
    def set_text(self, run, text, font, size, color):
        run.text = text
        run.font.name = font
        run.font.size = size
        run.font.color.rgb = color
