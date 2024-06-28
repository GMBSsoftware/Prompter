from WordReader import WordReader
from PPTCreator import PPTCreator
from Setting import Pattern
from Setting import PPT_WORD
from Setting import Symbol
from Text import Sentence, Word, Paragraph
from Default import Default, JHI, JJS, HMH, LWD
from Default import JHI
import re, os, sys
from pptx.util import Pt
from pptx.dml.color import RGBColor
from docx import Document


class WordPrompterCreator:
    def __init__(self, person) -> None:
        self.titles = []
        self.max_line = PPT_WORD.max_line
        self.max_byte = PPT_WORD.max_byte_in_one_line
        self.slides = []
        self.texts = []
        self.text = None
        self.word_reader = WordReader()
        self.file_name = None
        self.is_new_slide = False
        if person == "JHI":
            self.person = JHI()
        elif person == "JJS":
            self.person = JJS()
        elif person == "HMH":
            self.person = HMH()
        elif person == "LWD":
            self.person = LWD()
        else:
            self.person = Default()

    def process_first(self, texts):
        """워드 문서 읽어와서 파일명, 제목 저장. 말씀 시작 부분 위치 저장. 기본 폰트 저장."""
        return_texts = []
        is_before_bible = True
        bible_font = ""
        start_index = 0
        i = -1
        for sentence in texts:
            i += 1
            # "본문" 단어 나오기 이전
            if is_before_bible:
                # "본문"
                if bool(re.search(Pattern.bible_guide, str(sentence))):
                    is_before_bible = False
                    continue
                # o 월 o 일 oo 말씀
                if bool(re.search(Pattern.file_name_word, str(sentence))):
                    self.file_name = re.search(
                        Pattern.file_name_word, str(sentence)
                    ).group()
                    return_texts.append(sentence)
                # 공백이 아닌 경우에만 주제에 추가
                elif not str(sentence) == "":
                    self.titles.append(str(sentence))
                    return_texts.append(sentence)
                # 공백이면
                else:
                    return_texts.append(sentence)
            else:
                # 공백일 때 건너뛰기
                if str(sentence) == "":
                    return_texts.append(sentence)
                # 성경 구절 정규식일 때
                elif bool(re.search(Pattern.bible, str(sentence))):
                    bible_font = sentence[0].font
                    continue
                # 성경 구절 문단을 나눠썼을 때는 폰트로 구별.
                elif sentence[0].font == bible_font:
                    continue
                # 말씀 시작 부분
                else:
                    self.word_font = sentence[0].font
                    # 말씀 부분일 때 바로 아래줄에서 append 해주니까 +1
                    start_index = i + 1
                    return_texts.append(sentence)
                    break
        return_texts.extend(texts[start_index:])
        return return_texts

    def process_person(self, texts, person):
        return_texts = []
        is_vedio = False
        is_caption = False
        is_bible = False
        for sentence in texts:
            # 끝 표시
            if bool(re.search(Pattern.end, str(sentence))):
                if is_vedio:
                    is_vedio = False
                    # 사람마다 다르게. 제거하는 경우 return none 남기는 경우 텍스트 그대로 반환
                    result = person.process_end_vedio(sentence)
                    if result is None:
                        continue
                elif is_caption:
                    is_caption = False
                elif is_bible:
                    is_bible = False

            # 영상 내용인 경우 결과에 추가 x = 제거
            if is_vedio:
                continue
            # 영상 표시
            if bool(re.search(Pattern.vedio, str(sentence))):
                # 사람마다 is_vedio 다르게
                is_vedio = person.process_vedio(sentence)
                return_texts.append(sentence)
            else:
                return_texts.append(sentence)
        return return_texts

    def resource_path(self, relative_path):
        """Get absolute path to resource, works for dev and for PyInstaller"""
        try:
            # PyInstaller creates a temp folder and stores path in _MEIPASS
            base_path = sys._MEIPASS
        except AttributeError:
            base_path = os.path.abspath(".")
        return os.path.join(base_path, relative_path)

    def make_prompter(self, file):
        doc = Document(file)
        self.ppt = PPTCreator(
            PPT_WORD.back_color, self.resource_path("sample_word.pptx")
        )

        slide = self.ppt.add_new_slide()
        text_words = self.word_reader.convert(doc)

        texts = self.process_first(text_words)
        texts = self.process_person(texts, self.person)

        texts = self.process_line(texts)
        for paragraph in text_words:
            slide = self.max_process(paragraph, slide)

        # self.reLine()

        desktop_directory = os.path.join(os.path.expanduser("~"), "Desktop")
        self.ppt.prs.save(f"{desktop_directory}/{self.file_name}.pptx")

    # utf-8로 인코드 했을 때 텍스트의 바이트 구하는 메서드
    def length(self, text):
        if (
            isinstance(text, Sentence)
            or isinstance(text, Word)
            or isinstance(text, Paragraph)
        ):
            text = str(text)
        return len(text.encode("utf-8"))

    def resize_text(self, text_last, text_first):
        """길이 다른 2 텍스트 입력 받아서 합친 후 절반 잘라서 리턴 (길이 맞춰서 리턴)"""
        words = []
        for word in text_first:
            words.append(word)
        for word in text_last:
            words.append(word)
        half_length = len(words) // 2
        sentences = []
        sentences.append(Sentence(words[:half_length]))
        sentences.append(Sentence(words[half_length:]))

        return sentences

    # 텍스트 or 텍스트 리스트들의 길이가 초과했는지 체크
    def check_over_length(self, text_or_texts, max_byte):
        # 텍스트가 한 개면 리스트화해서 반복
        if not isinstance(text_or_texts, list):
            # 텍스트가 이미 \n로 나뉘었으면 리스트로 분리
            if "\n" in text_or_texts:
                text_or_texts = text_or_texts.split("\n")
            else:
                text_or_texts = [text_or_texts]
        # 텍스트들 각각 체크
        for i in text_or_texts:
            if self.length(i) > max_byte:
                return True
        return False

    # 현재 슬라이드의 줄 수와 입력할 텍스트의 줄 수를 합친 값이 최대 줄 수 초과하는지 체크
    def check_over_line(self, text, slide=None):
        # 현재 슬라이드에 있는 줄 수에 변수로 넘긴 text를 합쳤을 때 최대 줄 수 넘으면 true
        slide_text = ""

        # 슬라이드에서 한 줄씩 가져와서 반복
        if slide:
            for paragraph in slide.shapes.title.text_frame.paragraphs:
                # 공백이면 문단 나눠진 거니 \n 추가
                if str(paragraph) == "":
                    slide_text += "\n"
                # 글자 있으면 추가
                else:
                    slide_text += str(paragraph) + "\n"
            slide_line = slide_text.strip().count("\n") + 1
        else:
            slide_line = 0

        text_line = 0

        if isinstance(text, Paragraph):
            if str(text) == "":
                text_line = 0
            else:
                text_line = len(text.sentences)
        elif isinstance(text, Sentence):
            text_line = 1
        elif isinstance(text, str):
            text_line = text.count("\n") + 1

        if slide_line + text_line > self.max_line:
            return True
        return False

    # 최대 줄 수 넘을 때 문단 나누는 메서드.
    def split_over_line(self, paragraph, max_line):
        num_sentences = len(paragraph.sentences)
        if num_sentences <= max_line * 2:
            mid_point = num_sentences // 2
            first_half = Paragraph(paragraph[:mid_point])
            second_half = Paragraph(paragraph[mid_point:])
            return [first_half, second_half]
        else:
            paragraphs = []
            start = 0
            while start < num_sentences:
                end = min(start + max_line, num_sentences)
                paragraphs.append(Paragraph(paragraph[start:end]))
                start = end
            return paragraphs

    def process_line(self, texts):
        """한 줄 넘는 문장 분리"""
        return_texts = []
        for text in texts:
            if self.check_over_length(str(text), self.max_byte):
                # 한 줄 초과 시
                if not self.check_over_length(
                    self.join_comma_ideal(text, self.max_byte), self.max_byte
                ):
                    # 컴마로 나눴을 때 이상적으로 나뉘면
                    return_texts.extend(self.join_comma_ideal(text, self.max_byte))
                else:
                    return_texts.extend(self.join(text, self.max_byte))
            else:
                return_texts.append(text)
        return return_texts

    def process_slide(self, texts):
        """최대 줄 수 넘으면 분리"""

    # 최대 글자, 최대 줄 수 넘는지 체크해서 넘으면 나누는 프로세스
    def max_process(self, paragraph, slide):
        # 최대 글자 초과시 분리, 재조합 프로세스 실행
        if self.check_over_length(str(paragraph), self.max_byte):
            # 컴마로 나눌 때 이상적으로 잘 나눠져서 길이 안 넘으면
            if not self.check_over_length(
                self.join_comma_ideal(str(paragraph), self.max_byte),
                self.max_byte,
            ):
                paragraph = self.join_comma_ideal(paragraph, self.max_byte)
            else:
                # 그냥 길이 맞춰서 이어붙이는 메서드
                paragraph = self.join(paragraph, self.max_byte)

        if self.check_new_slide(paragraph):
            self.texts.append(self.text)
            if isinstance(paragraph, Sentence):
                self.text = paragraph
            elif isinstance(paragraph, Paragraph):
                self.text = paragraph[0]
        else:
            self.text = paragraph

        """if self.check_over_line(paragraph):
            paragraphs = self.split_over_line(paragraph, self.max_line)
            for i in paragraphs:
                self.slides.append(slide)
                slide = self.ppt.add_new_slide()
                self.write_on_slide(slide, i)
            return slide"""

        # 기존 슬라이드 줄 수 + 현재 텍스트의 줄수가 최대 줄 수 초과
        if self.check_over_line(paragraph, slide):

            self.slides.append(slide)
            slide = self.ppt.add_new_slide()
            self.write_on_slide(slide, paragraph)
            return slide

        # 최대 줄 수 미만이라 이어 붙이기
        else:
            self.write_on_slide(slide, paragraph)
            return slide

    def check_new_slide(self, text):
        if isinstance(text, Sentence):
            check_symbol_text = str(text)
        elif isinstance(text, Paragraph):
            check_symbol_text = str(text[0])

        # 중요 기호들은 새 슬라이드에서 시작.
        # 이미 새 슬라이드면
        if self.is_new_slide:
            self.is_new_slide = False
            return False
        # 중요 기호 있거나 빈 슬라이드면 새 슬라이드에서 시작.
        elif (
            any(symbol in check_symbol_text for symbol in Symbol.symbol_important)
            or check_symbol_text == ""
        ):
            self.is_new_slide = True
            return True
        else:
            return False

    def join_comma_ideal(self, text, max_byte):
        """컴마를 기준으로 나눴을 때 이상적으로 나눠지면 나눠서 반환, 아니면 그대로 반환"""
        max_byte = self.max_byte
        comma_index = -1
        while True:
            for i in range(len(text)):
                if "," in text[i]:
                    comma_index = i
                text1 = text[: comma_index + 1]
                text2 = text[comma_index + 1 :]

    # 컴마를 기준으로 나눴을 때 이상적으로 나눠지면 나눠서 반환, 아니면 그대로 반환하는 메서드
    def join_comma_ideal(self, text, max_byte):
        comma_index = -1
        # text = str(text)
        while True:
            # 현재 콤마의 인덱스를 찾습니다.
            comma_index = text.find(",", comma_index + 1)
            # 더 이상 콤마가 없으면 루프를 종료합니다.
            if comma_index == -1:
                break
            # 텍스트를 콤마 위치를 기준으로 두 부분으로 나눕니다.
            text1 = text[: comma_index + 1]
            text2 = text[comma_index + 1 :]

            # 각 부분이 최대 바이트 수를 초과하지 않는지 확인합니다.
            if not self.check_over_length(
                text1, max_byte
            ) and not self.check_over_length(text2, max_byte):
                # 두 부분의 길이 차이가 서로 2배 이상 나지 않는지 확인합니다.
                if abs(self.length(text1) - self.length(text2)) <= min(
                    self.length(text1), self.length(text2)
                ):
                    return [Sentence(text1), Sentence(text2)]

        # 적절한 분할이 없으면 원래 텍스트를 반환합니다.
        return [Sentence(text)]

    def join(self, words, max_byte):
        """max_byte 길이 찰 때까지 쭉 이어붙이는 메서드\n
        들어오는 값은 Word 클래스"""
        return_sentences = []
        sentence = Sentence()
        max_byte = self.max_byte
        for word in words:
            if (self.length(sentence) + self.length(word)) < max_byte:
                # max_byte 이하일 때 쭉 이어붙이기
                sentence.add_word(word)
            else:
                # max_byte 넘어가서 반환할 배열에 추가 후 다시 반복
                return_sentences.append(sentence)
                sentence = Sentence()
                sentence.add_word(word)
        return_sentences.append(sentence)
        if len(return_sentences) > 1:
            return_sentences.extend(
                self.resize_text(return_sentences.pop(), return_sentences.pop())
            )
            # 마지막거 2개 가져와서 길이 맞추는거
        return return_sentences

    """# max_byte 길이 찰 때까지 쭉 이어붙이는 메서드
    def join(self, texts, max_byte):
        result = Sentence()
        return_paragraph = Paragraph()
        for word in texts:
            # max_byte 이하일 때 쭉 이어붙이기
            if (self.length(result) + self.length(word)) < max_byte:
                result.add_word(word)
            # max_byte 넘어가서 반환할 배열에 추가 후 다시 반복
            else:
                return_paragraph.add_sentence(result)
                result = Sentence()
                result.add_word(word)
        return_paragraph.add_sentence(result)
        if len(return_paragraph.sentences) > 1:
            last_text = Paragraph()
            last_text.add_sentence(return_paragraph.get_and_remove(-2))
            last_text.add_sentence(return_paragraph.get_and_remove(-1))
            # if len(return_paragraph.sentences) > 2:
            # return_paragraph = Paragraph(return_paragraph[:-2])
            splitted_text = self.split_text_half(last_text)
            for i in splitted_text:
                if isinstance(i, Sentence):
                    return_paragraph.add_sentence(i)
        return return_paragraph"""

    # ppt에 텍스트 쓰기
    def write_on_slide(self, slide, paragraph):
        title_shape = slide.shapes.title
        title_text_frame = title_shape.text_frame
        ppt_paragraph = title_text_frame.paragraphs[-1]  # 마지막 단락 선택
        if isinstance(paragraph, Paragraph):
            for sentence in paragraph:
                for word in sentence:
                    self.add_text(ppt_paragraph, word)
                run = ppt_paragraph.add_run()
                run.text = "\n"

        elif isinstance(paragraph, Sentence):
            for word in paragraph:
                self.add_text(ppt_paragraph, word)
            run = ppt_paragraph.add_run()
            run.text = "\n"

    # ppt 슬라이드에 글, 폰트, 색상, 굵은 글씨, 밑줄 적용해서 쓰기
    def add_text(self, ppt_paragraph, word):
        slide_run = ppt_paragraph.add_run()
        slide_run.text = word.text + " "
        slide_run.font.name = PPT_WORD.font
        slide_run.font.size = Pt(PPT_WORD.size)
        # 색상 없으면 기본 색상
        if word.color is None:
            slide_run.font.color.rgb = PPT_WORD.default_color
        # 색상 있으면 그 색상 그대로
        else:
            # 워드에서 받은 객체랑 ppt에 쓸 객체가 서로 달라서 직접 변환 시켜줘야됨
            color = str(word.color)
            slide_run.font.color.rgb = RGBColor(
                int(color[0:2], 16),
                int(color[2:4], 16),
                int(color[4:6], 16),
            )
        # 굵은 글씨
        if not word.bold is None:
            slide_run.font.bold = True

        # 밑줄
        if not word.underline is None:
            slide_run.font.underline = True

    def reLine(self):
        for i in range(len(self.slides)):
            print(self.slides[i].shapes.title.text)
            lines_num = self.count_lines(self.slides[i].shapes.title.text)
            print("줄 수 : ", lines_num)
            print("=================================")
            if lines_num == 1 and i > 0:
                tmp = self.slides[i - 1].shapes.title.text
                self.slides[i - 1].shapes.title.text = self.slides[i].shapes.title.text
                self.slides[i].shapes.title.text = tmp

    # 맨 위, 맨 아래 공백 제외 글자 줄 수
    def count_lines(self, text):
        # 텍스트를 줄 단위로 분리
        lines = text.splitlines()
        # 맨 앞과 맨 뒤의 공백 줄 제거
        while lines and lines[0].strip() == "":
            lines.pop(0)
        while lines and lines[-1].strip() == "":
            lines.pop(-1)
        # 남은 줄의 개수를 반환
        return len(lines)
